import fitz  # PyMuPDF
from pptx import Presentation
from typing import Optional
import os

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extracts all text from a PDF file."""
    doc_text = []
    try:
        with fitz.open(pdf_path) as doc:
            for page in doc:
                doc_text.append(page.get_text())
        return "\n".join(doc_text)
    except Exception as e:
        raise RuntimeError(f"Error extracting text from PDF: {str(e)}")

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extracts text from all slides in a PowerPoint file."""
    doc_text = []
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    doc_text.append(shape.text)
        return "\n".join(doc_text)
    except Exception as e:
        raise RuntimeError(f"Error extracting text from PPTX: {str(e)}")

def extract_text(file_path: str) -> str:
    """Generic text extractor that routes based on file extension."""
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
